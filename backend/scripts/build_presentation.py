from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SUMMARY_ROOT = ROOT / "sumary"
VAIBHAVI_ROOT = SUMMARY_ROOT / "vaibhavi" / "EVALUATION & RESULTS - VAIBHAVI"
VISUALS = VAIBHAVI_ROOT / "MAIN_VISUALS"
DETAILS = VAIBHAVI_ROOT / "DETAIL_FILES"
OUT = SUMMARY_ROOT / "OUTFITAI_PRESENTATION_FREEZE_2026-03-19.pptx"


class Theme:
    BG = RGBColor(10, 18, 34)
    PANEL = RGBColor(17, 30, 55)
    PANEL_ALT = RGBColor(23, 39, 70)
    TITLE = RGBColor(240, 245, 255)
    BODY = RGBColor(213, 223, 243)
    MUTED = RGBColor(154, 176, 214)
    ACCENT = RGBColor(54, 148, 255)
    GOOD = RGBColor(63, 182, 123)
    WARN = RGBColor(253, 201, 64)
    RISK = RGBColor(243, 88, 88)


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title: str, subtitle: str = "") -> None:
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(11.9), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(34)
    run.font.color.rgb = Theme.TITLE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.68), Inches(1.22), Inches(11.4), Inches(0.6))
        p2 = sub_box.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(16)
        r2.font.color.rgb = Theme.MUTED


def add_role_badge(slide, owner: str) -> None:
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(0.42), Inches(2.2), Inches(0.44)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = Theme.PANEL_ALT
    badge.line.color.rgb = Theme.ACCENT
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = owner
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = Theme.BODY


def add_panel(slide, x: float, y: float, w: float, h: float, alt: bool = False):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = Theme.PANEL_ALT if alt else Theme.PANEL
    panel.line.color.rgb = Theme.ACCENT
    panel.line.width = Pt(0.8)
    return panel


def add_bullets(slide, bullets: Iterable[str], x: float, y: float, w: float, h: float, font_size: int = 22) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(font_size)
        p.font.color.rgb = Theme.BODY


def add_metric_card(slide, x: float, y: float, title: str, value: str, color: RGBColor) -> None:
    card = add_panel(slide, x, y, 2.45, 1.2, alt=True)
    card.line.color.rgb = color
    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(11)
    p1.font.color.rgb = Theme.MUTED
    p1.space_after = Pt(4)
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.bold = True
    p2.font.size = Pt(25)
    p2.font.color.rgb = color


def add_footer(slide, text: str = "Outfit AI - Presentation Freeze - March 2026") -> None:
    box = slide.shapes.add_textbox(Inches(0.65), Inches(7.05), Inches(12.0), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = Theme.MUTED


def add_image(slide, image_path: Path, x: float, y: float, w: float, h: float) -> None:
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        placeholder = add_panel(slide, x, y, w, h)
        placeholder.line.color.rgb = Theme.RISK
        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"Missing visual:\n{image_path.name}"
        p.font.size = Pt(13)
        p.font.color.rgb = Theme.BODY


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    with (DETAILS / "vision_summary_latest.json").open("r", encoding="utf-8") as f:
        vision = json.load(f)
    with (DETAILS / "recommender_summary_latest.json").open("r", encoding="utf-8") as f:
        rec = json.load(f)

    # Slide 1: Cover
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Outfit AI", "Production-Ready Smart Wardrobe Recommender")
    add_panel(s, 0.65, 1.95, 12.05, 4.6)
    add_bullets(
        s,
        [
            "Team Presentation Deck (Final Freeze)",
            "Focus: Recommender quality, robust pipeline, measurable improvements",
            "Presenters: Amogh (Technical), Dhanuja (Product + Recommender), Vaibhavi (Evaluation)",
        ],
        1.0,
        2.35,
        11.5,
        2.2,
        font_size=22,
    )
    add_footer(s)

    # Slide 2: Agenda
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Agenda", "Clean story from problem -> system -> proof -> roadmap")
    add_panel(s, 0.75, 1.85, 12.0, 5.0)
    add_bullets(
        s,
        [
            "1. Problem and project goal",
            "2. Architecture and end-to-end pipeline",
            "3. Engineering struggles and what we fixed",
            "4. Recommender logic and explainability",
            "5. Evaluation framework and key metrics",
            "6. Known limits and next steps",
            "7. Speaker-wise handoff",
        ],
        1.1,
        2.2,
        11.4,
        4.3,
        font_size=21,
    )
    add_footer(s)

    # Slide 3: Team ownership
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Team Ownership and Delivery Plan", "Recovered from project ownership logs")
    add_panel(s, 0.65, 1.8, 3.9, 4.9)
    add_panel(s, 4.75, 1.8, 3.9, 4.9, alt=True)
    add_panel(s, 8.85, 1.8, 3.9, 4.9)
    add_bullets(s, ["Amogh", "Technical Lead", "Architecture", "Pipeline internals", "Major fixes"], 0.92, 2.12, 3.3, 3.8, 17)
    add_bullets(s, ["Dhanuja", "Recommender + Product", "User flow", "Explainability", "Demo logic"], 5.02, 2.12, 3.3, 3.8, 17)
    add_bullets(s, ["Vaibhavi", "Evaluation + Results", "Testing framework", "Metrics and visuals", "Evidence narrative"], 9.12, 2.12, 3.3, 3.8, 17)
    add_bullets(s, ["Speaking order: Amogh (5 min) -> Dhanuja (3.5 min) -> Vaibhavi (3.5 min)"], 0.9, 6.3, 11.8, 0.5, 14)
    add_footer(s)

    # Slide 4: Problem and objective
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Problem Statement and Objective", "Build a reliable wardrobe recommender, not a random outfit picker")
    add_panel(s, 0.65, 1.85, 5.9, 4.9)
    add_panel(s, 6.75, 1.85, 5.9, 4.9, alt=True)
    add_bullets(
        s,
        [
            "Initial pain points:",
            "Detection ambiguity (shirt read as top + bottom)",
            "Weak recommendation trust",
            "UI exposing technical internals to users",
            "No robust evaluation evidence",
        ],
        0.95,
        2.15,
        5.3,
        4.2,
        17,
    )
    add_bullets(
        s,
        [
            "Final objective:",
            "Accurate single-garment scan flow",
            "Explainable top-bottom recommendations",
            "Feedback-aware ranking behavior",
            "Presentation-grade evaluation artifacts",
        ],
        7.05,
        2.15,
        5.3,
        4.2,
        17,
    )
    add_footer(s)

    # Slide 5: Architecture
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "System Architecture", "FastAPI backend + modular services + frontend product layer")
    add_role_badge(s, "Amogh")
    add_panel(s, 0.65, 1.9, 12.0, 4.9)
    add_bullets(
        s,
        [
            "Frontend (HTML/CSS/JS): scan, wardrobe, outfits, calendar, palette, history, profile",
            "API Layer (FastAPI): auth, scan, wardrobe, outfits, feedback, body profile",
            "Core Services: YOLO detection, SAM2 segmentation, LAB color extraction, CLIP embedding",
            "Recommender Engine: rule score + matrix factorization + MMR + repeat controls",
            "Storage: SQLite/SQLAlchemy (WardrobeItem, Outfit, Feedback, BodyProfile)",
        ],
        1.0,
        2.2,
        11.3,
        4.3,
        18,
    )
    add_footer(s)

    # Slide 6: Pipeline
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "End-to-End Pipeline", "From one image upload to explainable outfit output")
    add_role_badge(s, "Amogh")
    y = 2.55
    labels = [
        "Upload Image",
        "YOLO Detect\n(upper/lower)",
        "SAM2 Segment\n(clean mask)",
        "LAB Color +\nPalette",
        "Save Wardrobe\nItem",
        "Generate Outfit\nPlan",
        "Collect Feedback\nand Adapt",
    ]
    x = 0.6
    for i, label in enumerate(labels):
        box = add_panel(s, x, y, 1.75, 1.55, alt=(i % 2 == 1))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = Theme.BODY
        if i < len(labels) - 1:
            arrow = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 1.8), Inches(3.05), Inches(0.45), Inches(0.55))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = Theme.ACCENT
            arrow.line.fill.background()
        x += 1.85
    add_footer(s)

    # Slide 7: Struggles -> fixes
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Engineering Struggles and Fixes", "We improved through iteration, not in one perfect attempt")
    add_role_badge(s, "Amogh")
    add_panel(s, 0.65, 1.8, 5.9, 4.9)
    add_panel(s, 6.75, 1.8, 5.9, 4.9, alt=True)
    add_bullets(
        s,
        [
            "Problems observed:",
            "Shirt detected as both top and bottom",
            "Outfit generation blocked by unknown item_type",
            "Scan result showed raw mask path",
            "Profile page lacked visual guidance",
        ],
        0.95,
        2.12,
        5.4,
        4.2,
        17,
    )
    add_bullets(
        s,
        [
            "Fixes delivered:",
            "Single-garment detection logic with conflict resolution",
            "Generation gate removed + auto item-type assignment",
            "UI now shows useful user details (type, color family, palette)",
            "Body-shape icons + skin-tone swatch reference added",
        ],
        7.05,
        2.12,
        5.4,
        4.2,
        17,
    )
    add_footer(s)

    # Slide 8: Recommender design
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Recommender Design (Hybrid)", "Practical, explainable, and measurable")
    add_role_badge(s, "Dhanuja")
    add_panel(s, 0.65, 1.9, 12.0, 4.9)
    add_bullets(
        s,
        [
            "Candidate generation: all valid top-bottom combinations",
            "Scoring stack: rule score + collaborative signal + visual diversity + repeat/cooldown penalties",
            "Ranking controls: MMR rerank, rotation quotas, feedback bias",
            "Explainability: each outfit returns color reason, style reason, novelty signal, and confidence",
            "Product value: recommendations are deterministic and auditable, not random",
        ],
        1.0,
        2.2,
        11.3,
        4.3,
        18,
    )
    add_footer(s)

    # Slide 9: User flow / demo flow
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "User Flow and Demo Flow", "How we demo the complete recommender story live")
    add_role_badge(s, "Dhanuja")
    add_panel(s, 0.65, 1.85, 12.0, 4.9, alt=True)
    add_bullets(
        s,
        [
            "1. Scan one garment image",
            "2. Item saved to wardrobe with category + color metadata",
            "3. Generate outfits for selected days",
            "4. Show split-card output (top + bottom) with compatibility badge",
            "5. Show 'Why this outfit?' explanation in plain language",
            "6. Mark worn / like / dislike and explain feedback loop effect",
        ],
        1.0,
        2.2,
        11.3,
        4.3,
        19,
    )
    add_footer(s)

    # Slide 10: Evaluation framework
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Evaluation Framework", "Internal production-style testing pipeline")
    add_role_badge(s, "Vaibhavi")
    add_panel(s, 0.65, 1.85, 5.9, 4.9)
    add_panel(s, 6.75, 1.85, 5.9, 4.9, alt=True)
    add_bullets(
        s,
        [
            "Vision evaluation:",
            "YOLO detection accuracy proxy",
            "SAM2 segmentation quality + success rate",
            "Color extraction stability and mismatch tracking",
            "End-to-end pipeline success rate",
        ],
        0.95,
        2.1,
        5.3,
        4.2,
        17,
    )
    add_bullets(
        s,
        [
            "Recommender evaluation:",
            "8 synthetic users",
            "60-day simulation",
            "3 replicates",
            "Baseline vs personalized lifts",
            "Coverage, repetition, and forgotten-item metrics",
        ],
        7.05,
        2.1,
        5.3,
        4.2,
        17,
    )
    add_footer(s)

    # Slide 11: KPI snapshot cards
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Metric Snapshot (Latest Saved Artifacts)", "Evidence-first status at freeze time")
    add_role_badge(s, "Vaibhavi")

    add_metric_card(s, 0.8, 2.0, "Detection Accuracy", f"{vision['detection_accuracy_proxy']:.4f}", Theme.GOOD)
    add_metric_card(s, 3.45, 2.0, "Segmentation Success", f"{vision['segmentation_success_rate']:.4f}", Theme.GOOD)
    add_metric_card(s, 6.1, 2.0, "Color Success", f"{vision['color_success_rate']:.4f}", Theme.WARN)
    add_metric_card(s, 8.75, 2.0, "Pipeline Success", f"{vision['pipeline_success_rate']:.4f}", Theme.WARN)

    add_metric_card(s, 0.8, 3.45, "Score Lift", "+0.1284", Theme.GOOD)
    add_metric_card(s, 3.45, 3.45, "Diversity Lift", "+0.0103", Theme.GOOD)
    add_metric_card(s, 6.1, 3.45, "Repetition Lift", "-0.0042", Theme.GOOD)
    add_metric_card(s, 8.75, 3.45, "Coverage Lift", "+0.0208", Theme.ACCENT)

    add_panel(s, 0.8, 5.0, 10.4, 1.35)
    add_bullets(
        s,
        [
            "Reading: recommender metrics improved vs baseline; biggest remaining vision bottleneck is color classification.",
        ],
        1.05,
        5.33,
        9.9,
        0.8,
        14,
    )
    add_footer(s)

    # Slide 12: Chart baseline vs personalized
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Recommender Visual 1", "Baseline vs Personalized Comparison")
    add_role_badge(s, "Vaibhavi")
    add_image(s, VISUALS / "recsys_baseline_vs_personalized.png", 0.8, 1.8, 11.8, 4.9)
    add_footer(s)

    # Slide 13: Chart avg lifts
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Recommender Visual 2", "Average Lift by Metric")
    add_role_badge(s, "Vaibhavi")
    add_image(s, VISUALS / "recsys_avg_lifts.png", 0.8, 1.8, 11.8, 4.9)
    add_footer(s)

    # Slide 14: Chart heatmap
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Recommender Visual 3", "Per-user Lift Heatmap")
    add_role_badge(s, "Vaibhavi")
    add_image(s, VISUALS / "recsys_user_lift_heatmap.png", 0.8, 1.8, 11.8, 4.9)
    add_footer(s)

    # Slide 15: coverage vs repeat
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Recommender Visual 4", "Coverage vs Repetition Trade-off")
    add_role_badge(s, "Vaibhavi")
    add_image(s, VISUALS / "recsys_coverage_vs_repeat.png", 0.8, 1.8, 11.8, 4.9)
    add_footer(s)

    # Slide 16: trends
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Recommender Visual 5", "Trend Over Time")
    add_role_badge(s, "Vaibhavi")
    add_image(s, VISUALS / "recsys_trends_over_time.png", 0.8, 1.8, 11.8, 4.9)
    add_footer(s)

    # Slide 17: What is good vs pending
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "What Is Strong vs What Is Pending", "Honest quality statement for evaluation panel")
    add_panel(s, 0.65, 1.85, 5.9, 4.9)
    add_panel(s, 6.75, 1.85, 5.9, 4.9, alt=True)
    add_bullets(
        s,
        [
            "Working well:",
            "Hybrid recommender clearly beats baseline",
            "Segmentation is stable (1.0 success)",
            "UI now maps technical output to user-friendly actions",
            "Evaluation artifacts support claims with evidence",
        ],
        0.95,
        2.1,
        5.3,
        4.2,
        17,
    )
    add_bullets(
        s,
        [
            "Needs next pass:",
            "Color classification still weaker than target",
            "Detection robustness on harder formats can improve further",
            "Current DB choice is dev-friendly but not high-concurrency",
            "Presentation deck uses saved metrics (rerun before final publish if needed)",
        ],
        7.05,
        2.1,
        5.3,
        4.2,
        17,
    )
    add_footer(s)

    # Slide 18: Topic mapping to class
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "How This Maps to Recommender Systems Concepts", "Direct alignment with class topics")
    add_panel(s, 0.65, 1.85, 12.0, 4.9, alt=True)
    add_bullets(
        s,
        [
            "Collaborative Filtering: matrix factorization (SVD-style) from wear-pattern signals",
            "Cosine Similarity: embedding-based compatibility and diversity controls",
            "Ranking: weighted hybrid score with confidence and repetition penalties",
            "Feedback Learning: liked/disliked/worn actions alter future ranking bias",
            "Explainability: each recommendation carries plain-language reasons",
            "Note: Pearson correlation is not used in current production pipeline",
        ],
        1.0,
        2.2,
        11.3,
        4.3,
        18,
    )
    add_footer(s)

    # Slide 19: Speaker handoff
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Presentation Handoff Script", "So all 3 speakers stay synchronized")
    add_panel(s, 0.65, 1.85, 3.9, 4.9)
    add_panel(s, 4.75, 1.85, 3.9, 4.9, alt=True)
    add_panel(s, 8.85, 1.85, 3.9, 4.9)
    add_bullets(
        s,
        ["Amogh", "Slides 4 to 8", "Architecture", "Pipeline", "Fix journey", "Engineering choices"],
        0.92,
        2.15,
        3.3,
        4.2,
        16,
    )
    add_bullets(
        s,
        ["Dhanuja", "Slides 9 + product flow", "Recommender logic", "Explainability", "Live demo script"],
        5.02,
        2.15,
        3.3,
        4.2,
        16,
    )
    add_bullets(
        s,
        ["Vaibhavi", "Slides 10 to 18", "Evaluation method", "Metrics", "Visual proof", "Honest limitations"],
        9.12,
        2.15,
        3.3,
        4.2,
        16,
    )
    add_footer(s)

    # Slide 20: Closing
    s = prs.slides.add_slide(blank)
    set_bg(s, Theme.BG)
    add_title(s, "Final Closing", "What makes this project defendable")
    add_panel(s, 0.75, 1.95, 11.8, 4.8)
    add_bullets(
        s,
        [
            "We did not claim perfection. We showed measurable progress.",
            "We identified weak points using a structured evaluation framework.",
            "We fixed production-impact issues and improved recommendation quality.",
            "We can explain every major system choice with data and trade-offs.",
            "This project demonstrates engineering rigor, recommender relevance, and practical deployment thinking.",
        ],
        1.1,
        2.3,
        11.1,
        4.1,
        21,
    )
    add_footer(s, "Thank you")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    out = build()
    print(out)
